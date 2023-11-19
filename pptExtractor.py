import os
import tkinter as tk
from tkinter import filedialog, simpledialog
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH

def get_slide_text(slide):
    slide_text = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            slide_text.append(shape.text)
    return "\n".join(slide_text)

def search_phrases_in_presentation(file_path, phrases):
    prs = Presentation(file_path)
    slides_content = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = get_slide_text(slide)
        if any(phrase.lower() in slide_text.lower() for phrase in phrases if phrase):
            slides_content.append((i, slide_text))
    return slides_content

def search_phrases_in_folder(folder_path, phrases):
    results = {}
    for id, filename in enumerate(sorted(os.listdir(folder_path)), start=1):
        if filename.endswith(".pptx"):
            file_path = os.path.join(folder_path, filename)
            slides_content = search_phrases_in_presentation(file_path, phrases)
            if slides_content:
                results[f"{id:03d}_{filename}"] = slides_content
    return results

def add_footer_with_page_number_placeholder(doc):
    for section in doc.sections:
        footer = section.footer
        paragraph = footer.paragraphs[0]
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = paragraph.add_run("Page [Page Number]")
        run.font.size = Pt(9)

def save_results_to_word(results, output_file, cover_text):
    doc = Document()
    doc.add_heading('PowerPoint Search Report', 0)
    doc.add_paragraph(cover_text)
    doc.add_page_break()

    for filename, slides in results.items():
        doc.add_heading(filename, level=1)
        for slide_number, slide_text in slides:
            doc.add_heading(f"Slide {slide_number}", level=2)
            doc.add_paragraph(slide_text)
            doc.add_page_break()

    add_footer_with_page_number_placeholder(doc)
    doc.save(output_file)

def run_search():
    folder_path = filedialog.askdirectory(title="Select Folder with PowerPoint files")
    if not folder_path:
        return

    phrase1 = simpledialog.askstring("Input", "Enter the first phrase to search for:")
    phrase2 = simpledialog.askstring("Input", "Enter the second phrase to search for (optional):")
    phrase3 = simpledialog.askstring("Input", "Enter the third phrase to search for (optional):")

    if not phrase1 and not phrase2 and not phrase3:
        tk.messagebox.showwarning("No Input", "At least one phrase must be entered.")
        return

    output_word_file = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word Document", "*.docx")])
    if not output_word_file:
        return

    cover_text = simpledialog.askstring("Cover Page", "Enter the cover page text:")

    phrases = [phrase1, phrase2, phrase3]
    results = search_phrases_in_folder(folder_path, phrases)
    if results:
        save_results_to_word(results, output_word_file, cover_text)
        tk.messagebox.showinfo("Completed", "The search is complete and the results have been saved.")
    else:
        tk.messagebox.showinfo("No Results", "No slides found with the specified phrases.")

root = tk.Tk()
root.title("PowerPoint Phrase Search")

run_button = tk.Button(root, text="Run Search", command=run_search)
run_button.pack(pady=20)

root.mainloop()
