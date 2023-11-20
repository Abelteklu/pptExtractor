import os
import tkinter as tk
from tkinter import filedialog, simpledialog
from pptx import Presentation
from pptx.util import Inches

def copy_slide_from_external_pres(source_presentation, slide_index, target_presentation):
    # Create a new slide in the target presentation with same layout as the source slide
    source_slide = source_presentation.slides[slide_index]
    layout = source_slide.slide_layout
    new_slide = target_presentation.slides.add_slide(layout)

    # Copy all shapes from the source slide to the new slide
    for shape in source_slide.shapes:
        el = shape.element
        new_el = new_slide.shapes._spTree.addnext(el)

    return new_slide

def search_phrases_in_presentation(file_path, phrases):
    prs = Presentation(file_path)
    matched_slides = []
    for i, slide in enumerate(prs.slides):
        slide_text = ' '.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        if any(phrase.lower() in slide_text.lower() for phrase in phrases if phrase):
            matched_slides.append(i)
    return matched_slides

def search_phrases_in_folder(folder_path, phrases):
    results = {}
    for id, filename in enumerate(sorted(os.listdir(folder_path)), start=1):
        if filename.endswith(".pptx"):
            file_path = os.path.join(folder_path, filename)
            matched_slides = search_phrases_in_presentation(file_path, phrases)
            if matched_slides:
                results[f"{id:03d}_{filename}"] = (file_path, matched_slides)
    return results

def save_results_to_pptx(results, output_folder):
    for file_id, (file_path, slide_indices) in results.items():
        source_presentation = Presentation(file_path)
        target_presentation = Presentation()

        for slide_index in slide_indices:
            copy_slide_from_external_pres(source_presentation, slide_index, target_presentation)
        
        target_file_path = os.path.join(output_folder, f"Extracted_{file_id}")
        target_presentation.save(target_file_path)

def run_search():
    folder_path = filedialog.askdirectory(title="Select Folder with PowerPoint files")
    if not folder_path:
        return

    phrase1 = simpledialog.askstring("Input", "Enter the first phrase to search for:")
    phrase2 = simpledialog.askstring("Input", "Enter the second phrase to search for (optional):")
    phrase3 = simpledialog.askstring("Input", "Enter the third phrase to search for (optional):")

    if not phrase1 and not phrase2 and not phrase3:
        tk.messagebox.showwarning("No Input", "At least one phrase must be entered.")
        return

    output_folder = filedialog.askdirectory(title="Select Folder to Save Extracted Slides")
    if not output_folder:
        return

    phrases = [phrase1, phrase2, phrase3]
    results = search_phrases_in_folder(folder_path, phrases)
    if results:
        save_results_to_pptx(results, output_folder)
        tk.messagebox.showinfo("Completed", "The search is complete and the results have been saved.")
    else:
        tk.messagebox.showinfo("No Results", "No slides found with the specified phrases.")

root = tk.Tk()
root.title("PowerPoint Phrase Search")

run_button = tk.Button(root, text="Run Search", command=run_search)
run_button.pack(pady=20)

root.mainloop()
